import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import streamlit as st

# Load environment variables from .env file
load_dotenv()

# Configure Google Generative AI
genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))

def get_text_from_pdf(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_from_excel(excel_docs):
    text = ""
    for excel in excel_docs:
        df = pd.read_excel(excel)
        text += df.to_string()
    return text

def get_text_from_word(word_docs):
    text = ""
    for word in word_docs:
        doc = Document(word)
        for para in doc.paragraphs:
            text += para.text
    return text

def get_text_from_ppt(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text

def get_text_from_txt(txt_docs):
    text = ""
    for txt in txt_docs:
        with open(txt, "r", encoding="utf-8") as file:
            text += file.read()
    return text

def get_text_from_csv(csv_docs):
    text = ""
    for csv in csv_docs:
        df = pd.read_csv(csv)
        text += df.to_string()
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not in
    the provided context, just say, "answer is not available in the context." Do not provide a wrong answer.

    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    # Set allow_dangerous_deserialization to True
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)

    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def main():
    st.set_page_config(page_title="Chat PDF")
    st.header("Chat with PDF using Gemini💁")

    user_question = st.text_input("Ask a Question from the Uploaded Files")

    if st.button("Search"):
        if user_question:
            with st.spinner("Generating response..."):
                response = user_input(user_question)
                st.write("Reply: ", response)
        else:
            st.error("Please enter a question.")

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
            if uploaded_files:
                with st.spinner("Processing..."):
                    raw_text = ""
                    pdf_docs = [file for file in uploaded_files if file.name.endswith('.pdf')]
                    excel_docs = [file for file in uploaded_files if file.name.endswith('.xlsx') or file.name.endswith('.xls')]
                    word_docs = [file for file in uploaded_files if file.name.endswith('.docx')]
                    ppt_docs = [file for file in uploaded_files if file.name.endswith('.pptx')]
                    txt_docs = [file for file in uploaded_files if file.name.endswith('.txt')]
                    csv_docs = [file for file in uploaded_files if file.name.endswith('.csv')]

                    if pdf_docs:
                        raw_text += get_text_from_pdf(pdf_docs)
                    if excel_docs:
                        raw_text += get_text_from_excel(excel_docs)
                    if word_docs:
                        raw_text += get_text_from_word(word_docs)
                    if ppt_docs:
                        raw_text += get_text_from_ppt(ppt_docs)
                    if txt_docs:
                        raw_text += get_text_from_txt(txt_docs)
                    if csv_docs:
                        raw_text += get_text_from_csv(csv_docs)

                    if raw_text:
                        text_chunks = get_text_chunks(raw_text)
                        get_vector_store(text_chunks)
                        st.success("Done")
                    else:
                        st.error("No supported files found.")
            else:
                st.error("Please upload at least one file.")

if __name__ == "__main__":
    main()
